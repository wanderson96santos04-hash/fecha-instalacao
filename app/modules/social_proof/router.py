from __future__ import annotations

from datetime import datetime, timezone
from io import BytesIO
from typing import Optional

from fastapi import APIRouter, Request, Depends, Form
from fastapi.responses import HTMLResponse, RedirectResponse, Response
from fastapi.templating import Jinja2Templates
from sqlalchemy.orm import Session

from app.core.deps import pop_flashes, get_user_id_from_request
from app.db.session import SessionLocal
from app.models.user import User


router = APIRouter(prefix="/app/social-proof", tags=["social-proof"])
templates = Jinja2Templates(directory="app/templates")


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def get_current_user(request: Request, db: Session) -> User:
    uid = get_user_id_from_request(request)
    if not uid:
        # main.py já redireciona 401 -> /login
        raise Exception("Unauthorized")
    user = db.query(User).filter(User.id == int(uid)).first()
    if not user:
        raise Exception("Unauthorized")
    return user


def _format_money(valor: str) -> str:
    v = (valor or "").strip()
    if not v:
        return ""
    return v.replace("R$", "").strip()


def _build_testimonial(servico: str, valor: str, cidade: str, detalhe: str) -> str:
    servico = (servico or "").strip()
    cidade = (cidade or "").strip()
    detalhe = (detalhe or "").strip()

    val = _format_money(valor)
    prefix = "✅ Fechamos mais um serviço!"

    parts = [prefix, ""]
    parts.append(f"📌 Serviço: {servico}")
    if val:
        parts.append(f"💰 Valor: R$ {val}")
    if cidade:
        parts.append(f"📍 Cidade: {cidade}")
    if detalhe:
        parts.append(f"⭐ {detalhe}")

    parts.append("")
    parts.append("Se você também quer orçamento rápido e serviço bem feito, me chama no WhatsApp. 🔥")

    return "\n".join(parts).strip() + "\n"


def _ctx_base(request: Request, user: User, flashes, form: dict, result: str):
    """
    Importante:
    - O template social_proof.html usa a variável `result`
    - Mantemos `result_text` também por compatibilidade, mas o principal é `result`
    """
    now = datetime.now(timezone.utc)
    return {
        "request": request,
        "user": user,
        "flashes": flashes,
        "now": now,
        "result": result,          # ✅ o template usa isso
        "result_text": result,     # (opcional) compat
        "form": form,
    }


@router.get("", response_class=HTMLResponse)
def social_proof_page(request: Request, db: Session = Depends(get_db)):
    user = get_current_user(request, db)
    flashes = pop_flashes(request)

    return templates.TemplateResponse(
        "social_proof/social_proof.html",
        _ctx_base(
            request=request,
            user=user,
            flashes=flashes,
            form={"servico": "", "valor": "", "cidade": "", "detalhe": ""},
            result="",
        ),
    )


@router.post("/generate", response_class=HTMLResponse)
def social_proof_generate(
    request: Request,
    servico: str = Form(...),
    valor: str = Form(""),
    cidade: str = Form(""),
    detalhe: str = Form(""),
    db: Session = Depends(get_db),
):
    user = get_current_user(request, db)
    flashes = pop_flashes(request)

    text = _build_testimonial(servico=servico, valor=valor, cidade=cidade, detalhe=detalhe)

    return templates.TemplateResponse(
        "social_proof/social_proof.html",
        _ctx_base(
            request=request,
            user=user,
            flashes=flashes,
            form={"servico": servico, "valor": valor, "cidade": cidade, "detalhe": detalhe},
            result=text,
        ),
    )


def _require_pro(user: User) -> Optional[RedirectResponse]:
    if not getattr(user, "is_pro", False):
        return RedirectResponse(url="/app/upgrade", status_code=303)
    return None


@router.post("/pdf")
def social_proof_pdf(
    request: Request,
    text: str = Form(""),
    db: Session = Depends(get_db),
):
    user = get_current_user(request, db)
    redir = _require_pro(user)
    if redir:
        return redir

    if not text.strip():
        return RedirectResponse(url="/app/social-proof", status_code=303)

    # Importa aqui (não quebra o app no startup)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception:
        return Response(
            content="PDF export indisponível no servidor (dependência não instalada).",
            status_code=503,
            media_type="text/plain",
        )

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4  # noqa: F841

    c.setTitle("Prova Social - FECHA INSTALAÇÃO")
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, height - 60, "FECHA INSTALAÇÃO — Prova Social")

    c.setFont("Helvetica", 10)
    c.drawString(40, height - 80, datetime.now(timezone.utc).strftime("%d/%m/%Y %H:%M UTC"))

    y = height - 120
    c.setFont("Helvetica", 11)

    for line in text.splitlines():
        if y < 60:
            c.showPage()
            y = height - 60
            c.setFont("Helvetica", 11)
        c.drawString(40, y, line[:1200])
        y -= 16

    c.showPage()
    c.save()

    pdf_bytes = buf.getvalue()
    buf.close()

    filename = "prova-social.pdf"
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/ppt")
def social_proof_ppt(
    request: Request,
    text: str = Form(""),
    db: Session = Depends(get_db),
):
    user = get_current_user(request, db)
    redir = _require_pro(user)
    if redir:
        return redir

    if not text.strip():
        return RedirectResponse(url="/app/social-proof", status_code=303)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return Response(
            content="PPT export indisponível no servidor (dependência não instalada).",
            status_code=503,
            media_type="text/plain",
        )

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # título
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FECHA INSTALAÇÃO — Prova Social"
    run.font.size = Pt(28)
    run.font.bold = True

    # conteúdo
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(5.2))
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.clear()

    for i, line in enumerate(text.splitlines()):
        if i == 0:
            tf2.text = line
            tf2.paragraphs[0].font.size = Pt(18)
            tf2.paragraphs[0].font.bold = True
        else:
            p = tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    out = BytesIO()
    prs.save(out)
    ppt_bytes = out.getvalue()
    out.close()

    filename = "prova-social.pptx"
    return Response(
        content=ppt_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )