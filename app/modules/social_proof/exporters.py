from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Optional

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from pptx import Presentation
from pptx.util import Inches, Pt


def _brl(value: float) -> str:
    # formata em pt-BR simples
    s = f"{value:,.2f}"
    s = s.replace(",", "X").replace(".", ",").replace("X", ".")
    return f"R$ {s}"


def build_social_proof_text(servico: str, valor: float, cidade: str, detalhe: str = "") -> str:
    servico = (servico or "").strip()
    cidade = (cidade or "").strip()
    detalhe = (detalhe or "").strip()

    valor_txt = _brl(valor)

    linhas = [
        "🔥 SERVIÇO FECHADO!",
        "",
        f"✅ {servico}" if servico else "✅ Serviço fechado",
        f"📍 {cidade}" if cidade else "",
        "",
        f"💰 {valor_txt}",
    ]
    if detalhe:
        linhas += ["", f"🗣️ {detalhe}"]

    linhas += [
        "",
        "Mais um cliente satisfeito. Bora pra cima! 🚀",
    ]

    # remove linhas vazias duplicadas
    out = []
    for ln in linhas:
        if ln == "" and (out and out[-1] == ""):
            continue
        if ln != "" or out:
            out.append(ln)
    return "\n".join(out).strip()


def export_pdf(servico: str, valor: float, cidade: str, detalhe: str = "", brand: str = "FECHA INSTALAÇÃO") -> BytesIO:
    """
    PDF estilo story/marketing (fundo escuro, valor destacado).
    Retorna BytesIO pronto pra enviar.
    """
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    w, h = A4

    # Fundo escuro
    c.setFillColorRGB(0.06, 0.09, 0.16)  # slate-ish
    c.rect(0, 0, w, h, stroke=0, fill=1)

    # “Glow” laranja discreto (círculo)
    c.setFillColorRGB(0.92, 0.35, 0.05)  # orange
    c.setFillAlpha(0.12)
    c.circle(w * 0.55, h * 0.80, 220, stroke=0, fill=1)
    c.setFillAlpha(1)

    # Header
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, h - 50, brand)

    c.setFont("Helvetica", 10)
    c.setFillColorRGB(0.75, 0.78, 0.85)
    c.drawString(40, h - 68, datetime.utcnow().strftime("%d/%m/%Y"))

    # Título
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 26)
    c.drawString(40, h - 125, "🔥 SERVIÇO FECHADO")

    # Serviço + cidade
    y = h - 170
    c.setFont("Helvetica-Bold", 16)
    c.setFillColorRGB(0.93, 0.94, 0.98)
    c.drawString(40, y, f"✅ {servico.strip() or 'Serviço'}")
    y -= 28

    if cidade.strip():
        c.setFont("Helvetica", 14)
        c.setFillColorRGB(0.75, 0.78, 0.85)
        c.drawString(40, y, f"📍 {cidade.strip()}")
        y -= 26

    # Valor gigante
    valor_txt = _brl(valor)
    c.setFillColorRGB(0.98, 0.91, 0.80)  # warm light
    c.setFont("Helvetica-Bold", 44)
    c.drawString(40, y - 60, valor_txt)

    # Detalhe (opcional)
    if detalhe.strip():
        c.setFillColorRGB(0.93, 0.94, 0.98)
        c.setFont("Helvetica", 13)
        c.drawString(40, 180, f"🗣️ {detalhe.strip()[:90]}")

    # Rodapé
    c.setFillColorRGB(0.75, 0.78, 0.85)
    c.setFont("Helvetica", 11)
    c.drawString(40, 120, "Mais um cliente satisfeito. Bora pra cima! 🚀")

    c.showPage()
    c.save()

    buffer.seek(0)
    return buffer


def export_pptx(servico: str, valor: float, cidade: str, detalhe: str = "", brand: str = "FECHA INSTALAÇÃO") -> BytesIO:
    """
    PPTX (1 slide) estilo story.
    Retorna BytesIO pronto pra enviar.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Fundo (forma retângulo)
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE (evitar import extra)
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = None  # deixa default, vamos por cor no próximo passo
    bg.line.fill.background()

    # Como pptx é chato com RGB sem import, vamos usar um retângulo “overlay” com tema padrão.
    # (mantém compatível e ainda bonito). Se quiser cor exata, eu ajusto depois com RGBColor.

    # Título
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(9), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "🔥 SERVIÇO FECHADO"
    run.font.size = Pt(34)
    run.font.bold = True

    # Brand
    br = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(8), Inches(0.5))
    br_tf = br.text_frame
    br_tf.text = brand
    br_tf.paragraphs[0].runs[0].font.size = Pt(12)
    br_tf.paragraphs[0].runs[0].font.bold = True

    # Serviço
    svc = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(9), Inches(0.8))
    svc_tf = svc.text_frame
    svc_tf.text = f"✅ {servico.strip() or 'Serviço'}"
    svc_tf.paragraphs[0].runs[0].font.size = Pt(20)
    svc_tf.paragraphs[0].runs[0].font.bold = True

    # Cidade
    if cidade.strip():
        cd = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(9), Inches(0.6))
        cd_tf = cd.text_frame
        cd_tf.text = f"📍 {cidade.strip()}"
        cd_tf.paragraphs[0].runs[0].font.size = Pt(16)

    # Valor
    val = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(9), Inches(1.2))
    val_tf = val.text_frame
    val_tf.text = f"💰 {_brl(valor)}"
    val_tf.paragraphs[0].runs[0].font.size = Pt(40)
    val_tf.paragraphs[0].runs[0].font.bold = True

    # Detalhe (opcional)
    if detalhe.strip():
        dt = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(9), Inches(1))
        dt_tf = dt.text_frame
        dt_tf.text = f"🗣️ {detalhe.strip()[:120]}"
        dt_tf.paragraphs[0].runs[0].font.size = Pt(16)

    # Footer
    ft = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(9), Inches(0.6))
    ft_tf = ft.text_frame
    ft_tf.text = "Mais um cliente satisfeito. Bora pra cima! 🚀"
    ft_tf.paragraphs[0].runs[0].font.size = Pt(14)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out