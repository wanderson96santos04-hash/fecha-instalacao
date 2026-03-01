from __future__ import annotations

import os
import re
import io
from datetime import datetime, timezone, timedelta
from typing import Dict, List

from fastapi import APIRouter, Request, Form, HTTPException
from fastapi.responses import HTMLResponse, RedirectResponse, Response
from fastapi.templating import Jinja2Templates
from sqlalchemy import select, desc

from app.core.deps import get_user_id_from_request, redirect, pop_flashes
from app.db.session import SessionLocal
from app.models.user import User
from app.models.budget import Budget
from app.services.budget_service import can_create_budget, create_budget, FREE_LIMIT_TOTAL_BUDGETS
from app.services.whatsapp import build_budget_message, whatsapp_link, followup_message
from app.services.followup import can_followup

router = APIRouter(prefix="/app")
templates = Jinja2Templates(directory="app/templates")


def _require_user(request: Request) -> int:
    uid_raw = get_user_id_from_request(request)
    if not uid_raw:
        raise HTTPException(status_code=401)
    try:
        return int(uid_raw)
    except (TypeError, ValueError):
        raise HTTPException(status_code=401)


def _parse_brl_value(value: str) -> float:
    if not value:
        return 0.0
    s = str(value).strip()
    if not s:
        return 0.0

    s = re.sub(r"[^0-9\.,]", "", s)
    if not s:
        return 0.0

    if "," in s:
        s = s.replace(".", "").replace(",", ".")
        try:
            return float(s)
        except Exception:
            return 0.0

    if "." in s:
        parts = s.split(".")
        if len(parts) == 2 and len(parts[1]) in (1, 2):
            try:
                return float(s)
            except Exception:
                return 0.0
        s = s.replace(".", "")
        try:
            return float(s)
        except Exception:
            return 0.0

    try:
        return float(s)
    except Exception:
        return 0.0


def _money_brl(v: float) -> str:
    s = f"{v:,.2f}"
    s = s.replace(",", "X").replace(".", ",").replace("X", ".")
    return f"R$ {s}"


def _month_window_utc(now: datetime) -> tuple[datetime, datetime]:
    start = datetime(now.year, now.month, 1, tzinfo=timezone.utc)
    if now.month == 12:
        end = datetime(now.year + 1, 1, 1, tzinfo=timezone.utc)
    else:
        end = datetime(now.year, now.month + 1, 1, tzinfo=timezone.utc)
    return start, end


# ✅ FIX: aceitar status PT/EN para o painel do mês (sem mudar o banco)
def _norm_status(s: str) -> str:
    s = (s or "").strip().lower()
    if s in ("won", "fechado", "fechado (mês)", "close", "closed"):
        return "won"
    if s in ("lost", "perdido", "perdido (mês)"):
        return "lost"
    if s in ("awaiting", "aguardando", "pendente", "aguardando (mês)"):
        return "awaiting"
    return s


@router.get("", response_class=HTMLResponse)
def dashboard(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    now = datetime.now(timezone.utc)
    month_start, month_end = _month_window_utc(now)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

        budgets = list(
            db.scalars(
                select(Budget)
                .where(Budget.user_id == uid)
                .order_by(desc(Budget.created_at), desc(Budget.id))
            ).all()
        )

        total = len(budgets)
        remaining = None
        if not user.is_pro:
            remaining = max(0, FREE_LIMIT_TOTAL_BUDGETS - total)

        month_budgets = [
            b for b in budgets
            if b.created_at
            and (b.created_at.replace(tzinfo=timezone.utc) >= month_start)
            and (b.created_at.replace(tzinfo=timezone.utc) < month_end)
        ]

        # ✅ FIX: painel agora conta mesmo se tiver "fechado/perdido/aguardando"
        won = [b for b in month_budgets if _norm_status(b.status or "") == "won"]
        lost = [b for b in month_budgets if _norm_status(b.status or "") == "lost"]
        awaiting = [b for b in month_budgets if _norm_status(b.status or "") == "awaiting"]

        won_value = sum(_parse_brl_value(b.value or "") for b in won)
        lost_value = sum(_parse_brl_value(b.value or "") for b in lost)

        total_month = len(month_budgets)
        conversion_pct = (len(won) / total_month * 100.0) if total_month > 0 else 0.0

        metrics = {
            "month_won_value": _money_brl(won_value),
            "month_won_count": len(won),
            "month_lost_value": _money_brl(lost_value),
            "month_lost_count": len(lost),
            "month_conversion_pct": f"{conversion_pct:.0f}%",
            "month_total_count": len(month_budgets),
            "month_awaiting": len(awaiting),
            "month_awaiting_count": len(awaiting),
        }

    return templates.TemplateResponse(
        "dashboard.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "budgets": budgets,
            "total": total,
            "remaining": remaining,
            "can_followup": can_followup,
            "metrics": metrics,
        },
    )


@router.get("/acquisition", response_class=HTMLResponse)
def acquisition_page(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

        # Aquisição é PRO — FREE vai para /app/upgrade
        if not user.is_pro:
            return redirect(
                "/app/upgrade",
                kind="error",
                message="Esse módulo é exclusivo para usuários Premium.",
            )

    return templates.TemplateResponse(
        "acquisition.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "now": datetime.now(timezone.utc),
            "messages": [],
            "form": {"nicho": "", "cidade": "", "servico": "", "mode": "media"},
            "mode": "media",
        },
    )


def _generate_messages(nicho: str, cidade: str, servico: str, mode: str) -> List[str]:
    nicho = (nicho or "").strip()
    cidade = (cidade or "").strip()
    servico = (servico or "").strip()
    mode = (mode or "media").strip().lower()

    base = {
        "curta": [
            f"Olá! Trabalho com {servico} em {cidade}. Posso te enviar uma estimativa gratuita?",
            f"Oi! Atendo clientes de {nicho} em {cidade}. Quer saber quanto custaria o {servico} no seu caso?",
            f"Olá! Faço {servico} em {cidade}. Posso te passar um valor aproximado sem compromisso.",
            f"Boa tarde! Você já considerou {servico}? Posso te explicar rapidamente como funciona em {cidade}.",
            f"Olá! Muitos clientes em {cidade} estão procurando {servico}. Quer que eu te envie uma estimativa?",
            f"Oi! Trabalho com {servico} na região de {cidade}. Posso te mandar uma simulação gratuita.",
            f"Olá! Posso te passar uma orientação rápida sobre {servico} em {cidade}. Sem compromisso.",
            f"Oi! Atendo projetos de {servico} em {cidade}. Quer ver quanto ficaria no seu caso?",
            f"Olá! Posso te mostrar quanto você economizaria com {servico} em {cidade}.",
            f"Boa tarde! Faço {servico} em {cidade}. Quer receber uma estimativa gratuita?",
        ],
        "media": [
            f"Olá! Tudo bem? Trabalho com {servico} em {cidade}. Posso te enviar uma estimativa gratuita baseada no seu perfil?",
            f"Oi! Atendo clientes que buscam {servico} em {cidade}. Posso te passar uma simulação rápida sem compromisso.",
            f"Olá! Faço projetos de {servico} em {cidade}. Muitos clientes conseguem ótimo custo-benefício. Quer que eu te envie uma estimativa?",
            f"Boa tarde! Posso te explicar rapidamente como funciona o {servico} e quanto ficaria em média no seu caso em {cidade}.",
            f"Olá! Trabalho com {servico} na região de {cidade}. Posso te enviar uma previsão de investimento e retorno.",
            f"Oi! Atendo projetos de {nicho} em {cidade}. Posso te mandar uma orientação inicial e estimativa gratuita.",
            f"Olá! Muitos clientes em {cidade} estão procurando {servico}. Posso te mostrar quanto ficaria no seu caso.",
            f"Boa tarde! Posso te enviar uma simulação personalizada de {servico} para sua realidade em {cidade}.",
            f"Olá! Faço atendimento especializado em {servico}. Quer receber uma estimativa sem compromisso?",
            f"Oi! Posso te enviar uma projeção realista de custo e benefício do {servico} em {cidade}.",
        ],
        "agressiva": [
            f"Olá! Trabalho com {servico} em {cidade}. Posso te enviar uma estimativa gratuita hoje mesmo.",
            f"Oi! Muitos clientes em {cidade} estão iniciando {servico}. Posso te mostrar quanto ficaria no seu caso.",
            f"Olá! Posso te enviar uma simulação completa de {servico} com valores atualizados.",
            f"Boa tarde! Faço {servico} em {cidade}. Quer receber uma estimativa sem compromisso?",
            f"Olá! Posso te mostrar quanto você economizaria com {servico}.",
            f"Oi! Atendo clientes em {cidade}. Posso te enviar uma previsão de investimento.",
            f"Olá! Trabalho com instalação profissional de {servico}. Quer ver uma estimativa?",
            f"Boa tarde! Posso te enviar uma simulação gratuita e personalizada.",
            f"Olá! Posso te explicar rapidamente os valores do {servico} em {cidade}.",
            f"Oi! Quer receber uma estimativa gratuita e sem compromisso?",
        ],
    }

    return base.get(mode, base["media"])


@router.post("/acquisition/generate", response_class=HTMLResponse)
def acquisition_generate(
    request: Request,
    nicho: str = Form(...),
    cidade: str = Form(...),
    servico: str = Form(...),
    mode: str = Form("media"),
):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

        # POST também precisa bloquear FREE
        if not user.is_pro:
            return redirect(
                "/app/upgrade",
                kind="error",
                message="Esse módulo é exclusivo para usuários Premium.",
            )

    messages = _generate_messages(nicho=nicho, cidade=cidade, servico=servico, mode=mode)

    return templates.TemplateResponse(
        "acquisition.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "now": datetime.now(timezone.utc),
            "messages": messages,
            "form": {"nicho": nicho, "cidade": cidade, "servico": servico, "mode": mode},
            "mode": mode,
        },
    )


# =========================
# ROTAS DO MENU
# =========================
# ⚠️ IMPORTANTE:
# NÃO DEFINIR /onboarding AQUI.
# O onboarding verdadeiro está em app/modules/onboarding/router.py (e templates do módulo).
# Se definir aqui, quebra (TemplateNotFound) e ainda conflita rota.


@router.get("/invite", response_class=HTMLResponse)
def invite_page(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    invite_link = "https://fecha-instalacao.onrender.com/signup"
    copy_count = 0
    click_count = 0
    share_text = "Vem testar o sistema: https://fecha-instalacao.onrender.com/signup"

    return templates.TemplateResponse(
        "invite/invite.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "invite_link": invite_link,
            "copy_count": copy_count,
            "click_count": click_count,
            "share_text": share_text,
        },
    )


@router.get("/cases", response_class=HTMLResponse)
def cases_page(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    items: List[Dict] = []

    return templates.TemplateResponse(
        "cases/cases.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "items": items,
            "now": datetime.now(timezone.utc),
        },
    )


@router.get("/cases/admin", response_class=HTMLResponse)
def cases_admin_list(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    items: List[Dict] = []

    return templates.TemplateResponse(
        "cases/admin_list.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "items": items,
            "now": datetime.now(timezone.utc),
        },
    )


@router.get("/cases/admin/new", response_class=HTMLResponse)
def cases_admin_new(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    return templates.TemplateResponse(
        "cases/admin_new.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "now": datetime.now(timezone.utc),
        },
    )


@router.post("/cases/admin/new")
def cases_admin_new_post(request: Request):
    return RedirectResponse(url="/app/cases/admin", status_code=302)


@router.get("/cases/admin/edit/{item_id}", response_class=HTMLResponse)
def cases_admin_edit(request: Request, item_id: int):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    item = None

    return templates.TemplateResponse(
        "cases/admin_edit.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "item": item,
            "item_id": item_id,
            "now": datetime.now(timezone.utc),
        },
    )


@router.post("/cases/admin/edit/{item_id}")
def cases_admin_edit_post(request: Request, item_id: int):
    return RedirectResponse(url="/app/cases/admin", status_code=302)


@router.get("/cases/export", response_class=HTMLResponse)
def cases_export(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    return templates.TemplateResponse(
        "cases/export.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "now": datetime.now(timezone.utc),
        },
    )


@router.get("/social-proof", response_class=HTMLResponse)
def social_proof_page(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    return templates.TemplateResponse(
        "social_proof/social_proof.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "form": {"servico": "", "valor": "", "cidade": "", "detalhe": ""},
            "result": "",
            "now": datetime.now(timezone.utc),
        },
    )


@router.post("/social-proof/generate", response_class=HTMLResponse)
def social_proof_generate(
    request: Request,
    servico: str = Form(""),
    valor: str = Form(""),
    cidade: str = Form(""),
    detalhe: str = Form(""),
):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

    servico_s = (servico or "").strip()
    valor_s = (valor or "").strip()
    cidade_s = (cidade or "").strip()
    detalhe_s = (detalhe or "").strip()

    partes = []
    if servico_s:
        partes.append(f"Serviço fechado: {servico_s}")
    if valor_s:
        partes.append(f"Valor: {valor_s}")
    if cidade_s:
        partes.append(f"Cidade: {cidade_s}")
    if detalhe_s:
        partes.append(f"Detalhe: {detalhe_s}")

    if partes:
        result = "✅ Prova social pronta:\n" + " • ".join(partes)
    else:
        result = "Preencha o formulário para gerar a prova social."

    return templates.TemplateResponse(
        "social_proof/social_proof.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "form": {"servico": servico_s, "valor": valor_s, "cidade": cidade_s, "detalhe": detalhe_s},
            "result": result,
            "now": datetime.now(timezone.utc),
        },
    )


# =========================
# EXPORT (PDF / PPT) - CORREÇÃO REAL DO 405
# Aceita GET e POST e retorna o arquivo (sem redirect).
# =========================

def _sp_get_payload(
    request: Request,
    servico: str | None = None,
    valor: str | None = None,
    cidade: str | None = None,
    detalhe: str | None = None,
) -> dict:
    qp = request.query_params
    return {
        "servico": (servico if servico is not None else qp.get("servico", "")).strip(),
        "valor": (valor if valor is not None else qp.get("valor", "")).strip(),
        "cidade": (cidade if cidade is not None else qp.get("cidade", "")).strip(),
        "detalhe": (detalhe if detalhe is not None else qp.get("detalhe", "")).strip(),
    }


def _sp_text(payload: dict) -> str:
    partes = []
    if payload["servico"]:
        partes.append(f"Serviço fechado: {payload['servico']}")
    if payload["valor"]:
        partes.append(f"Valor: {payload['valor']}")
    if payload["cidade"]:
        partes.append(f"Cidade: {payload['cidade']}")
    if payload["detalhe"]:
        partes.append(f"Detalhe: {payload['detalhe']}")
    if not partes:
        return "Prova social (vazia). Preencha os campos antes de exportar."
    return "Prova social\n\n" + "\n".join(partes)


@router.get("/social-proof/pdf")
@router.post("/social-proof/pdf")
def social_proof_pdf(
    request: Request,
    servico: str = Form(""),
    valor: str = Form(""),
    cidade: str = Form(""),
    detalhe: str = Form(""),
):
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")
        if not user.is_pro:
            return redirect("/app/upgrade", kind="error", message="Exportação PDF é Premium.")

    payload = _sp_get_payload(request, servico, valor, cidade, detalhe)

    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
    except Exception:
        raise HTTPException(status_code=500, detail="Biblioteca de PDF não instalada (reportlab).")

    now = datetime.now(timezone.utc)
    filename = f'prova-social-{now.strftime("%Y%m%d-%H%M%S")}.pdf'

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4

    # Layout base
    margin = 54  # ~1.9cm
    card_x = margin
    card_y = margin + 36
    card_w = width - (margin * 2)
    card_h = height - (margin * 2) - 36

    # Fundo (branco) + "card" com borda suave
    c.setFillColor(colors.white)
    c.rect(0, 0, width, height, stroke=0, fill=1)

    c.setFillColor(colors.white)
    c.setStrokeColor(colors.HexColor("#E5E7EB"))  # cinza claro
    c.setLineWidth(1)
    c.roundRect(card_x, card_y, card_w, card_h, radius=14, stroke=1, fill=1)

    # Cabeçalho
    title_y = card_y + card_h - 54
    c.setFillColor(colors.HexColor("#0F172A"))  # slate-900
    c.setFont("Helvetica-Bold", 26)
    c.drawString(card_x + 28, title_y, "Prova Social")

    c.setFillColor(colors.HexColor("#475569"))  # slate-600
    c.setFont("Helvetica", 12)
    c.drawString(card_x + 28, title_y - 22, "Depoimento pronto para postar")

    # Linha separadora
    sep_y = title_y - 36
    c.setStrokeColor(colors.HexColor("#EEF2F7"))
    c.setLineWidth(1)
    c.line(card_x + 24, sep_y, card_x + card_w - 24, sep_y)

    # Conteúdo (labels + valores)
    y = sep_y - 34
    label_x = card_x + 28
    value_x = card_x + 150

    def draw_row(label: str, value: str, y_pos: float) -> float:
        if not value:
            return y_pos
        c.setFillColor(colors.HexColor("#334155"))  # slate-700
        c.setFont("Helvetica-Bold", 12)
        c.drawString(label_x, y_pos, label)

        c.setFillColor(colors.HexColor("#0F172A"))
        c.setFont("Helvetica", 12)
        # quebra simples se for longo
        max_w = (card_x + card_w - 28) - value_x
        words = value.split()
        line = ""
        lines: List[str] = []
        for w in words:
            test = (line + " " + w).strip()
            if c.stringWidth(test, "Helvetica", 12) <= max_w:
                line = test
            else:
                if line:
                    lines.append(line)
                line = w
        if line:
            lines.append(line)

        if not lines:
            lines = [""]

        c.drawString(value_x, y_pos, lines[0])
        y_pos -= 18
        for extra in lines[1:]:
            c.drawString(value_x, y_pos, extra)
            y_pos -= 18
        y_pos -= 4
        return y_pos

    y = draw_row("Serviço:", payload.get("servico", ""), y)
    y = draw_row("Valor:", payload.get("valor", ""), y)
    y = draw_row("Cidade:", payload.get("cidade", ""), y)
    y = draw_row("Detalhe:", payload.get("detalhe", ""), y)

    # Caso vazio: mensagem central
    if not (payload.get("servico") or payload.get("valor") or payload.get("cidade") or payload.get("detalhe")):
        c.setFillColor(colors.HexColor("#0F172A"))
        c.setFont("Helvetica-Bold", 16)
        c.drawString(card_x + 28, sep_y - 64, "Prova social vazia")
        c.setFillColor(colors.HexColor("#475569"))
        c.setFont("Helvetica", 12)
        c.drawString(card_x + 28, sep_y - 86, "Preencha os campos antes de exportar.")

    # Rodapé
    footer_y = card_y + 18
    c.setFillColor(colors.HexColor("#94A3B8"))  # slate-400
    c.setFont("Helvetica", 9)
    c.drawString(card_x + 24, footer_y, f"Gerado em {now.strftime('%d/%m/%Y %H:%M UTC')}")
    c.drawRightString(card_x + card_w - 24, footer_y, "Fecha Instalação")

    c.showPage()
    c.save()
    buf.seek(0)

    headers = {"Content-Disposition": f'attachment; filename="{filename}"'}
    return Response(content=buf.getvalue(), media_type="application/pdf", headers=headers)


@router.get("/social-proof/ppt")
@router.post("/social-proof/ppt")
def social_proof_ppt(
    request: Request,
    servico: str = Form(""),
    valor: str = Form(""),
    cidade: str = Form(""),
    detalhe: str = Form(""),
):
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")
        if not user.is_pro:
            return redirect("/app/upgrade", kind="error", message="Exportação PPT é Premium.")

    payload = _sp_get_payload(request, servico, valor, cidade, detalhe)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except Exception:
        raise HTTPException(status_code=500, detail="Biblioteca de PPT não instalada (python-pptx).")

    now = datetime.now(timezone.utc)
    filename = f'prova-social-{now.strftime("%Y%m%d-%H%M%S")}.pptx'

    prs = Presentation()
    # Layout em branco
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo (branco) + card
    bg = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.rectangle (evitar import extra)
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    card = slide.shapes.add_shape(
        1,
        Inches(0.7), Inches(0.7), Inches(11.9), Inches(6.1)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(229, 231, 235)
    card.line.width = Pt(1)

    # Título
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.0), Inches(10.8), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Prova Social"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(10.8), Inches(0.4))
    st = sub_box.text_frame
    st.clear()
    sp = st.paragraphs[0]
    sp.text = "Depoimento pronto para postar"
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(71, 85, 105)

    # Conteúdo
    body_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(10.8), Inches(3.5))
    bt = body_box.text_frame
    bt.clear()
    bt.word_wrap = True

    lines = []
    if payload["servico"]:
        lines.append(("Serviço", payload["servico"]))
    if payload["valor"]:
        lines.append(("Valor", payload["valor"]))
    if payload["cidade"]:
        lines.append(("Cidade", payload["cidade"]))
    if payload["detalhe"]:
        lines.append(("Detalhe", payload["detalhe"]))

    if not lines:
        p0 = bt.paragraphs[0]
        p0.text = "Sem dados — preencha os campos antes de exportar."
        p0.font.size = Pt(20)
        p0.font.color.rgb = RGBColor(71, 85, 105)
    else:
        # primeira linha
        k0, v0 = lines[0]
        p0 = bt.paragraphs[0]
        p0.text = f"{k0}: {v0}"
        p0.font.size = Pt(24)
        p0.font.color.rgb = RGBColor(15, 23, 42)

        for k, v in lines[1:]:
            p = bt.add_paragraph()
            p.text = f"{k}: {v}"
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(15, 23, 42)

    # Rodapé
    footer = slide.shapes.add_textbox(Inches(1.1), Inches(6.35), Inches(10.8), Inches(0.35))
    ft = footer.text_frame
    ft.clear()
    fp = ft.paragraphs[0]
    fp.text = f"Gerado em {now.strftime('%d/%m/%Y %H:%M UTC')}  •  Fecha Instalação"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(148, 163, 184)
    fp.alignment = PP_ALIGN.LEFT

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)

    headers = {"Content-Disposition": f'attachment; filename="{filename}"'}
    return Response(
        content=out.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


# ✅ NÃO coloque /app/upgrade aqui.
# Ela fica SOMENTE em app/routes/upgrade.py


# =========================
# NOVO ORÇAMENTO
# =========================

@router.get("/budgets/new", response_class=HTMLResponse)
def budgets_new_page(request: Request):
    flashes = pop_flashes(request)
    uid = _require_user(request)

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

        budgets = list(
            db.scalars(select(Budget).where(Budget.user_id == uid)).all()
        )

        total = len(budgets)
        remaining = None
        if not user.is_pro:
            remaining = max(0, FREE_LIMIT_TOTAL_BUDGETS - total)

    return templates.TemplateResponse(
        "budget_new.html",
        {
            "request": request,
            "flashes": flashes,
            "user": user,
            "total": total,
            "remaining": remaining,
            "can_create_budget": can_create_budget,
        },
    )


@router.post("/budgets/new")
def budgets_new_post(
    request: Request,
    client_name: str = Form(""),
    phone: str = Form(""),
    # ✅ FIX: create_budget exige service_type (mas aceitamos o campo antigo "service" também)
    service_type: str = Form("", alias="service_type"),
    service: str = Form(""),
    value: str = Form(""),
    payment_method: str = Form(""),
    notes: str = Form(""),
):
    uid = _require_user(request)

    client_name = (client_name or "").strip()
    phone = (phone or "").strip()
    service_type = (service_type or "").strip()
    service = (service or "").strip()
    value = (value or "").strip()
    payment_method = (payment_method or "").strip()
    notes = (notes or "").strip()

    # ✅ usa service_type se vier, senão usa service (compatível com template antigo)
    final_service_type = service_type or service

    with SessionLocal() as db:
        user = db.get(User, uid)
        if not user:
            return redirect("/login", kind="error", message="Faça login novamente.")

        if not can_create_budget(db, user):
            return redirect("/app/upgrade", kind="error", message="Você atingiu o limite do plano gratuito.")

        create_budget(
            db=db,
            user_id=uid,
            client_name=client_name,
            phone=phone,
            service_type=final_service_type,
            value=value,
            payment_method=payment_method,
            notes=notes,
        )

    return redirect("/app", kind="success", message="Orçamento criado com sucesso!")


# =========================
# WHATSAPP (ENVIO)
# =========================

@router.get("/budgets/{budget_id}/whatsapp")
def budgets_whatsapp(request: Request, budget_id: int):
    uid = _require_user(request)

    with SessionLocal() as db:
        budget = db.scalar(
            select(Budget).where(Budget.id == budget_id, Budget.user_id == uid)
        )
        if not budget:
            raise HTTPException(status_code=404, detail="Não encontrado")

        # ✅ FIX: NÃO chamar build_budget_message(budget) porque a função não aceita args
        # monta a mensagem aqui (só para essa rota, sem quebrar nada do resto)
        parts = []
        if (budget.service_type or "").strip():
            parts.append(f"Serviço: {budget.service_type}")
        if (budget.value or "").strip():
            parts.append(f"Valor: {budget.value}")
        if (budget.payment_method or "").strip():
            parts.append(f"Pagamento: {budget.payment_method}")
        if (budget.notes or "").strip():
            parts.append(f"Obs: {budget.notes}")

        msg = "Olá! Segue seu orçamento:\n" + "\n".join(parts) if parts else "Olá! Segue seu orçamento."

        phone = (budget.phone or "").strip()
        if not phone:
            return redirect("/app", kind="error", message="Esse orçamento não tem telefone cadastrado.")

        url = whatsapp_link(phone, msg)
        return RedirectResponse(url=url, status_code=302)


# =========================
# STATUS (MARCAR)
# =========================

@router.post("/budgets/{budget_id}/status")
def budgets_status_post(
    request: Request,
    budget_id: int,
    status: str = Form(""),
):
    uid = _require_user(request)
    status_s = (status or "").strip().lower()

    # ✅ FIX: aceita PT/EN mas grava padronizado
    mapping = {
        "awaiting": "awaiting",
        "aguardando": "awaiting",
        "pendente": "awaiting",

        "won": "won",
        "fechado": "won",

        "lost": "lost",
        "perdido": "lost",
    }

    status_db = mapping.get(status_s)
    if not status_db:
        return redirect("/app", kind="error", message="Status inválido.")

    with SessionLocal() as db:
        budget = db.scalar(
            select(Budget).where(Budget.id == budget_id, Budget.user_id == uid)
        )
        if not budget:
            raise HTTPException(status_code=404, detail="Não encontrado")

        budget.status = status_db
        db.add(budget)
        db.commit()

    return RedirectResponse(url="/app", status_code=303)